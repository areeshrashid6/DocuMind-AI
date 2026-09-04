"""
Turns an uploaded file's raw bytes into a list of LangChain Document objects
with useful metadata (page/sheet/slide, filename, source type).

Every loader function returns List[langchain_core.documents.Document].
"""

import io
import json
import os
import tempfile

import pandas as pd
from langchain_core.documents import Document

from src.ocr import image_bytes_to_text, pdf_is_scanned, ocr_pdf


def _tmp_path(file_bytes: bytes, suffix: str) -> str:
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(file_bytes)
        return tmp.name


def load_pdf(file_bytes: bytes, filename: str) -> list[Document]:
    from langchain_community.document_loaders import PyPDFLoader

    path = _tmp_path(file_bytes, ".pdf")
    try:
        if pdf_is_scanned(path):
            return ocr_pdf(path, filename)

        loader = PyPDFLoader(path)
        raw_docs = loader.load()
        docs = []
        for d in raw_docs:
            page = d.metadata.get("page", 0)
            docs.append(
                Document(
                    page_content=d.page_content,
                    metadata={
                        "source": filename,
                        "file_type": "pdf",
                        "page": page + 1,
                    },
                )
            )
        return docs
    finally:
        os.unlink(path)


def load_word(file_bytes: bytes, filename: str) -> list[Document]:
    from langchain_community.document_loaders import UnstructuredWordDocumentLoader

    ext = os.path.splitext(filename)[1].lower()
    path = _tmp_path(file_bytes, ext)
    try:
        loader = UnstructuredWordDocumentLoader(path)
        raw_docs = loader.load()
        return [
            Document(
                page_content=d.page_content,
                metadata={"source": filename, "file_type": "docx"},
            )
            for d in raw_docs
        ]
    finally:
        os.unlink(path)


def load_excel(file_bytes: bytes, filename: str) -> list[Document]:
    """
    Preserve sheet/row structure instead of dumping raw Unstructured text,
    so numeric questions ("highest revenue?") stay answerable.
    """
    ext = os.path.splitext(filename)[1].lower()
    engine = "openpyxl" if ext == ".xlsx" else None
    excel_file = pd.ExcelFile(io.BytesIO(file_bytes), engine=engine)

    docs = []
    for sheet_name in excel_file.sheet_names:
        df = excel_file.parse(sheet_name)
        df = df.dropna(how="all")
        if df.empty:
            continue

        columns = list(df.columns)
        # Batch rows into readable text chunks (~25 rows) to keep chunk sizes sane
        batch_size = 25
        for start in range(0, len(df), batch_size):
            batch = df.iloc[start : start + batch_size]
            lines = [f"Sheet: {sheet_name}", f"Columns: {', '.join(str(c) for c in columns)}", ""]
            for row_idx, row in batch.iterrows():
                pairs = ", ".join(f"{col}: {row[col]}" for col in columns)
                lines.append(f"Row {row_idx + 2}: {pairs}")
            docs.append(
                Document(
                    page_content="\n".join(lines),
                    metadata={
                        "source": filename,
                        "file_type": "xlsx",
                        "sheet": sheet_name,
                        "rows": f"{start + 2}-{start + len(batch) + 1}",
                    },
                )
            )
    return docs


def load_csv(file_bytes: bytes, filename: str, sep: str = ",") -> list[Document]:
    df = pd.read_csv(io.BytesIO(file_bytes), sep=sep)
    df = df.dropna(how="all")
    columns = list(df.columns)
    docs = []
    batch_size = 25
    for start in range(0, len(df), batch_size):
        batch = df.iloc[start : start + batch_size]
        lines = [f"Columns: {', '.join(str(c) for c in columns)}", ""]
        for row_idx, row in batch.iterrows():
            pairs = ", ".join(f"{col}: {row[col]}" for col in columns)
            lines.append(f"Row {row_idx + 2}: {pairs}")
        docs.append(
            Document(
                page_content="\n".join(lines),
                metadata={
                    "source": filename,
                    "file_type": "csv",
                    "rows": f"{start + 2}-{start + len(batch) + 1}",
                },
            )
        )
    return docs


def load_pptx(file_bytes: bytes, filename: str) -> list[Document]:
    from pptx import Presentation

    path = _tmp_path(file_bytes, ".pptx")
    try:
        prs = Presentation(path)
        docs = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if text.strip():
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text for cell in row.cells]
                        texts.append(" | ".join(cells))
            slide_text = "\n".join(texts).strip()
            if slide_text:
                docs.append(
                    Document(
                        page_content=slide_text,
                        metadata={"source": filename, "file_type": "pptx", "slide": i},
                    )
                )
        return docs
    finally:
        os.unlink(path)


def load_text(file_bytes: bytes, filename: str) -> list[Document]:
    text = file_bytes.decode("utf-8", errors="ignore")
    return [Document(page_content=text, metadata={"source": filename, "file_type": "text"})]


def load_html(file_bytes: bytes, filename: str) -> list[Document]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(file_bytes, "html.parser")
    text = soup.get_text(separator="\n")
    return [Document(page_content=text, metadata={"source": filename, "file_type": "html"})]


def load_xml(file_bytes: bytes, filename: str) -> list[Document]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(file_bytes, "xml")
    text = soup.get_text(separator="\n")
    return [Document(page_content=text, metadata={"source": filename, "file_type": "xml"})]


def load_json(file_bytes: bytes, filename: str) -> list[Document]:
    data = json.loads(file_bytes.decode("utf-8", errors="ignore"))
    text = json.dumps(data, indent=2, ensure_ascii=False)
    return [Document(page_content=text, metadata={"source": filename, "file_type": "json"})]


def load_image(file_bytes: bytes, filename: str) -> list[Document]:
    text = image_bytes_to_text(file_bytes)
    if not text.strip():
        text = "[No text detected by OCR in this image.]"
    return [
        Document(
            page_content=text,
            metadata={"source": filename, "file_type": "image", "ocr": True},
        )
    ]


def load_file(file_bytes: bytes, filename: str) -> list[Document]:
    """Route a file to the correct loader based on its extension."""
    ext = os.path.splitext(filename)[1].lower().lstrip(".")

    if ext == "pdf":
        return load_pdf(file_bytes, filename)
    if ext in ("doc", "docx", "rtf"):
        return load_word(file_bytes, filename)
    if ext in ("xls", "xlsx"):
        return load_excel(file_bytes, filename)
    if ext == "csv":
        return load_csv(file_bytes, filename, sep=",")
    if ext == "tsv":
        return load_csv(file_bytes, filename, sep="\t")
    if ext in ("ppt", "pptx"):
        return load_pptx(file_bytes, filename)
    if ext in ("txt", "md"):
        return load_text(file_bytes, filename)
    if ext in ("html", "htm"):
        return load_html(file_bytes, filename)
    if ext == "xml":
        return load_xml(file_bytes, filename)
    if ext == "json":
        return load_json(file_bytes, filename)
    if ext in ("png", "jpg", "jpeg", "webp", "tiff", "bmp"):
        return load_image(file_bytes, filename)

    raise ValueError(f"Unsupported file type: .{ext}")
